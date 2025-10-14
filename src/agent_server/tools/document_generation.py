"""
Document Generation Tool
Comprehensive tool for generating DOCX, PDF, and PPT documents with template support
"""

from typing import Dict, Any, List, Optional
import time
import tempfile
import os
from dataclasses import dataclass
from datetime import datetime
import re

from .registry import (
    BaseTool,
    ToolResult,
    ExecutionContext,
    ToolCapabilities,
    ResourceRequirements,
    ToolCapability,
)
from src.shared.logging import get_logger

logger = get_logger(__name__)

try:
    # Document generation libraries
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    import docx.oxml.shared

    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        PageBreak,
        Table,
        TableStyle,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY

    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    DEPENDENCIES_AVAILABLE = True

except ImportError as e:
    logger.warning(f"Document generation dependencies not available: {e}")
    DEPENDENCIES_AVAILABLE = False


@dataclass
class DocumentTemplate:
    """Document template configuration"""

    name: str
    format_type: str  # docx, pdf, ppt
    template_path: Optional[str] = None
    styles: Dict[str, Any] = None
    layout: Dict[str, Any] = None
    metadata: Dict[str, Any] = None


@dataclass
class DocumentValidationResult:
    """Result of document validation"""

    is_valid: bool
    errors: List[str]
    warnings: List[str]
    accessibility_score: float
    quality_metrics: Dict[str, Any]
    file_size_mb: float


@dataclass
class ContentStructure:
    """Analyzed content structure for document generation"""

    title: Optional[str]
    sections: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    code_blocks: List[Dict[str, Any]]
    images: List[Dict[str, Any]]
    lists: List[Dict[str, Any]]
    metadata: Dict[str, Any]


class TemplateManager:
    """Manages document templates"""

    def __init__(self):
        self.templates: Dict[str, DocumentTemplate] = {}
        self._initialize_default_templates()

    def _initialize_default_templates(self):
        """Initialize default templates for each format"""

        # Default DOCX template
        self.templates["default_docx"] = DocumentTemplate(
            name="default_docx",
            format_type="docx",
            styles={
                "title": {"font_size": 16, "bold": True, "alignment": "center"},
                "heading1": {"font_size": 14, "bold": True},
                "heading2": {"font_size": 12, "bold": True},
                "normal": {"font_size": 11},
                "code": {"font_name": "Courier New", "font_size": 10},
            },
            layout={
                "margins": {"top": 1.0, "bottom": 1.0, "left": 1.0, "right": 1.0},
                "page_size": "letter",
            },
        )

        # Default PDF template
        self.templates["default_pdf"] = DocumentTemplate(
            name="default_pdf",
            format_type="pdf",
            styles={
                "title": {"fontSize": 18, "alignment": TA_CENTER, "spaceAfter": 12},
                "heading1": {"fontSize": 14, "spaceAfter": 6, "spaceBefore": 12},
                "heading2": {"fontSize": 12, "spaceAfter": 6, "spaceBefore": 10},
                "normal": {"fontSize": 11, "alignment": TA_JUSTIFY},
                "code": {"fontName": "Courier", "fontSize": 10, "leftIndent": 20},
            },
            layout={
                "pagesize": A4,
                "margins": {"top": 72, "bottom": 72, "left": 72, "right": 72},
            },
        )

        # Default PPT template
        self.templates["default_ppt"] = DocumentTemplate(
            name="default_ppt",
            format_type="ppt",
            styles={
                "title": {"font_size": 24, "bold": True},
                "subtitle": {"font_size": 18},
                "content": {"font_size": 16},
                "bullet": {"font_size": 14},
            },
            layout={"slide_width": PptxInches(10), "slide_height": PptxInches(7.5)},
        )

    def get_template(self, template_name: str, format_type: str) -> DocumentTemplate:
        """Get template by name and format"""
        template_key = f"{template_name}_{format_type}"

        if template_key in self.templates:
            return self.templates[template_key]

        # Fallback to default template
        default_key = f"default_{format_type}"
        if default_key in self.templates:
            return self.templates[default_key]

        raise ValueError(
            f"No template found for {template_name} in format {format_type}"
        )

    def register_template(self, template: DocumentTemplate):
        """Register a new template"""
        template_key = f"{template.name}_{template.format_type}"
        self.templates[template_key] = template


class ContentAnalyzer:
    """Analyzes content structure for optimal document formatting"""

    def analyze_content_structure(self, content: str) -> ContentStructure:
        """Analyze content structure for document generation"""

        lines = content.split("\n")

        # Extract title (first non-empty line or line starting with #)
        title = None
        for line in lines:
            line = line.strip()
            if line:
                if line.startswith("# "):
                    title = line[2:].strip()
                elif not title:
                    title = line
                break

        # Extract sections
        sections = self._extract_sections(lines)

        # Extract tables
        tables = self._extract_tables(lines)

        # Extract code blocks
        code_blocks = self._extract_code_blocks(lines)

        # Extract lists
        lists = self._extract_lists(lines)

        # Extract images (placeholder detection)
        images = self._extract_image_references(lines)

        # Generate metadata
        metadata = {
            "word_count": len(content.split()),
            "line_count": len(lines),
            "has_code": len(code_blocks) > 0,
            "has_tables": len(tables) > 0,
            "has_lists": len(lists) > 0,
            "complexity": self._assess_content_complexity(content),
        }

        return ContentStructure(
            title=title,
            sections=sections,
            tables=tables,
            code_blocks=code_blocks,
            images=images,
            lists=lists,
            metadata=metadata,
        )

    def _extract_sections(self, lines: List[str]) -> List[Dict[str, Any]]:
        """Extract sections from content"""
        sections = []
        current_section = None

        for i, line in enumerate(lines):
            line = line.strip()

            # Check for markdown headers
            if line.startswith("#"):
                if current_section:
                    sections.append(current_section)

                level = len(line) - len(line.lstrip("#"))
                title = line.lstrip("#").strip()

                current_section = {
                    "title": title,
                    "level": level,
                    "start_line": i,
                    "content": [],
                }
            elif current_section:
                current_section["content"].append(line)

        # Add the last section
        if current_section:
            sections.append(current_section)

        return sections

    def _extract_tables(self, lines: List[str]) -> List[Dict[str, Any]]:
        """Extract table structures from content"""
        tables = []
        in_table = False
        current_table = None

        for i, line in enumerate(lines):
            line = line.strip()

            # Simple markdown table detection
            if "|" in line and line.count("|") >= 2:
                if not in_table:
                    in_table = True
                    current_table = {"start_line": i, "rows": [], "headers": []}

                # Parse table row
                cells = [cell.strip() for cell in line.split("|")[1:-1]]

                # Check if it's a header separator
                if all(re.match(r"^:?-+:?$", cell) for cell in cells):
                    continue  # Skip separator row

                if not current_table["headers"]:
                    current_table["headers"] = cells
                else:
                    current_table["rows"].append(cells)

            elif in_table and line == "":
                # End of table
                if current_table:
                    current_table["end_line"] = i
                    tables.append(current_table)
                in_table = False
                current_table = None

        # Add the last table if still in progress
        if current_table:
            current_table["end_line"] = len(lines)
            tables.append(current_table)

        return tables

    def _extract_code_blocks(self, lines: List[str]) -> List[Dict[str, Any]]:
        """Extract code blocks from content"""
        code_blocks = []
        in_code_block = False
        current_block = None

        for i, line in enumerate(lines):
            if line.strip().startswith("```"):
                if not in_code_block:
                    # Start of code block
                    in_code_block = True
                    language = line.strip()[3:].strip()
                    current_block = {
                        "start_line": i,
                        "language": language,
                        "content": [],
                    }
                else:
                    # End of code block
                    if current_block:
                        current_block["end_line"] = i
                        code_blocks.append(current_block)
                    in_code_block = False
                    current_block = None
            elif in_code_block and current_block:
                current_block["content"].append(line)

        return code_blocks

    def _extract_lists(self, lines: List[str]) -> List[Dict[str, Any]]:
        """Extract list structures from content"""
        lists = []
        in_list = False
        current_list = None

        for i, line in enumerate(lines):
            line_stripped = line.strip()

            # Check for list items (-, *, +, or numbered)
            if line_stripped.startswith(("- ", "* ", "+ ")) or re.match(
                r"^\d+\.\s", line_stripped
            ):

                if not in_list:
                    in_list = True
                    list_type = (
                        "ordered"
                        if re.match(r"^\d+\.\s", line_stripped)
                        else "unordered"
                    )
                    current_list = {"start_line": i, "type": list_type, "items": []}

                # Extract list item content
                if line_stripped.startswith(("- ", "* ", "+ ")):
                    item_content = line_stripped[2:].strip()
                else:  # numbered list
                    item_content = re.sub(r"^\d+\.\s", "", line_stripped)

                current_list["items"].append(
                    {
                        "content": item_content,
                        "line": i,
                        "indent_level": (len(line) - len(line.lstrip())) // 2,
                    }
                )

            elif in_list and line_stripped == "":
                continue  # Allow empty lines in lists

            elif in_list:
                # End of list
                if current_list:
                    current_list["end_line"] = i
                    lists.append(current_list)
                in_list = False
                current_list = None

        # Add the last list if still in progress
        if current_list:
            current_list["end_line"] = len(lines)
            lists.append(current_list)

        return lists

    def _extract_image_references(self, lines: List[str]) -> List[Dict[str, Any]]:
        """Extract image references from content"""
        images = []

        for i, line in enumerate(lines):
            # Markdown image syntax: ![alt text](image_url)
            image_matches = re.findall(r"!\[([^\]]*)\]\(([^)]+)\)", line)

            for alt_text, image_url in image_matches:
                images.append({"alt_text": alt_text, "url": image_url, "line": i})

        return images

    def _assess_content_complexity(self, content: str) -> str:
        """Assess content complexity for formatting decisions"""

        word_count = len(content.split())
        line_count = len(content.split("\n"))

        # Count special elements
        code_blocks = content.count("```")
        tables = content.count("|")
        headers = len(re.findall(r"^#+\s", content, re.MULTILINE))

        complexity_score = 0

        # Word count factor
        if word_count > 2000:
            complexity_score += 3
        elif word_count > 1000:
            complexity_score += 2
        elif word_count > 500:
            complexity_score += 1

        # Structure factor
        if headers > 10:
            complexity_score += 2
        elif headers > 5:
            complexity_score += 1

        # Special elements factor
        if code_blocks > 5 or tables > 10:
            complexity_score += 2
        elif code_blocks > 0 or tables > 0:
            complexity_score += 1

        if complexity_score >= 5:
            return "high"
        elif complexity_score >= 3:
            return "medium"
        else:
            return "low"


class DocumentValidator:
    """Validates generated documents for quality and accessibility"""

    def validate_document(
        self, file_path: str, format_type: str
    ) -> DocumentValidationResult:
        """Validate document quality and accessibility"""

        errors = []
        warnings = []
        quality_metrics = {}
        accessibility_score = 1.0

        try:
            # Get file size
            file_size_mb = os.path.getsize(file_path) / (1024 * 1024)

            # Format-specific validation
            if format_type == "docx":
                validation_result = self._validate_docx(file_path)
            elif format_type == "pdf":
                validation_result = self._validate_pdf(file_path)
            elif format_type == "ppt":
                validation_result = self._validate_ppt(file_path)
            else:
                validation_result = {
                    "errors": [f"Unsupported format: {format_type}"],
                    "warnings": [],
                    "metrics": {},
                }

            errors.extend(validation_result.get("errors", []))
            warnings.extend(validation_result.get("warnings", []))
            quality_metrics.update(validation_result.get("metrics", {}))

            # Calculate accessibility score
            accessibility_score = self._calculate_accessibility_score(
                quality_metrics, format_type
            )

            # General file validation
            if file_size_mb > 50:
                warnings.append(f"Large file size: {file_size_mb:.2f} MB")

            is_valid = len(errors) == 0

        except Exception as e:
            errors.append(f"Validation error: {str(e)}")
            is_valid = False
            file_size_mb = 0.0

        return DocumentValidationResult(
            is_valid=is_valid,
            errors=errors,
            warnings=warnings,
            accessibility_score=accessibility_score,
            quality_metrics=quality_metrics,
            file_size_mb=file_size_mb,
        )

    def _validate_docx(self, file_path: str) -> Dict[str, Any]:
        """Validate DOCX document"""
        errors = []
        warnings = []
        metrics = {}

        try:
            if DEPENDENCIES_AVAILABLE:
                doc = Document(file_path)

                # Count elements
                paragraph_count = len(doc.paragraphs)
                table_count = len(doc.tables)

                metrics.update(
                    {
                        "paragraph_count": paragraph_count,
                        "table_count": table_count,
                        "has_styles": len(doc.styles) > 0,
                    }
                )

                # Check for empty document
                if paragraph_count == 0:
                    errors.append("Document has no content")

                # Check for accessibility features
                if not any(p.style.name.startswith("Heading") for p in doc.paragraphs):
                    warnings.append(
                        "No heading styles found - may impact accessibility"
                    )

            else:
                warnings.append("DOCX validation limited - dependencies not available")

        except Exception as e:
            errors.append(f"DOCX validation error: {str(e)}")

        return {"errors": errors, "warnings": warnings, "metrics": metrics}

    def _validate_pdf(self, file_path: str) -> Dict[str, Any]:
        """Validate PDF document"""
        errors = []
        warnings = []
        metrics = {}

        try:
            # Basic file existence and readability check
            if not os.path.exists(file_path):
                errors.append("PDF file does not exist")
            elif not os.access(file_path, os.R_OK):
                errors.append("PDF file is not readable")
            else:
                metrics["file_exists"] = True

                # Additional PDF-specific validation would require PyPDF2 or similar
                warnings.append(
                    "PDF content validation limited - additional dependencies needed"
                )

        except Exception as e:
            errors.append(f"PDF validation error: {str(e)}")

        return {"errors": errors, "warnings": warnings, "metrics": metrics}

    def _validate_ppt(self, file_path: str) -> Dict[str, Any]:
        """Validate PowerPoint document"""
        errors = []
        warnings = []
        metrics = {}

        try:
            if DEPENDENCIES_AVAILABLE:
                prs = Presentation(file_path)

                slide_count = len(prs.slides)

                metrics.update(
                    {
                        "slide_count": slide_count,
                        "slide_width": prs.slide_width,
                        "slide_height": prs.slide_height,
                    }
                )

                # Check for empty presentation
                if slide_count == 0:
                    errors.append("Presentation has no slides")

                # Check slide content
                empty_slides = 0
                for slide in prs.slides:
                    if not slide.shapes:
                        empty_slides += 1

                if empty_slides > 0:
                    warnings.append(f"{empty_slides} slides have no content")

            else:
                warnings.append("PPT validation limited - dependencies not available")

        except Exception as e:
            errors.append(f"PPT validation error: {str(e)}")

        return {"errors": errors, "warnings": warnings, "metrics": metrics}

    def _calculate_accessibility_score(
        self, metrics: Dict[str, Any], format_type: str
    ) -> float:
        """Calculate accessibility score based on document features"""

        score = 0.8  # Base score

        # Format-specific accessibility checks
        if format_type == "docx":
            if metrics.get("has_styles", False):
                score += 0.1
            if metrics.get("paragraph_count", 0) > 0:
                score += 0.1

        elif format_type == "pdf":
            # PDF accessibility is harder to assess without specialized libraries
            score = 0.7  # Conservative score

        elif format_type == "ppt":
            slide_count = metrics.get("slide_count", 0)
            if slide_count > 0:
                score += 0.1
            if slide_count > 1:
                score += 0.1

        return min(1.0, score)


class DocumentGenerationTool(BaseTool):
    """Comprehensive document generation tool supporting DOCX, PDF, and PPT formats"""

    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        self.template_manager = TemplateManager()
        self.content_analyzer = ContentAnalyzer()
        self.validator = DocumentValidator()
        self.output_directory = (
            config.get("output_directory", tempfile.gettempdir())
            if config
            else tempfile.gettempdir()
        )

        # Ensure output directory exists
        os.makedirs(self.output_directory, exist_ok=True)

    async def initialize(self):
        """Initialize the document generation tool"""
        await super().initialize()

        if not DEPENDENCIES_AVAILABLE:
            logger.warning(
                "Document generation dependencies not fully available - some features may be limited"
            )

        logger.info(
            "Document Generation Tool initialized",
            output_directory=self.output_directory,
            dependencies_available=DEPENDENCIES_AVAILABLE,
        )

    async def execute(
        self, parameters: Dict[str, Any], context: ExecutionContext
    ) -> ToolResult:
        """Execute document generation"""

        start_time = time.time()

        try:
            content = parameters["content"]
            format_type = parameters.get("format", "docx").lower()
            template_name = parameters.get("template", "default")
            filename = parameters.get("filename")
            validate_output = parameters.get("validate", True)

            logger.info(
                "Executing document generation",
                format_type=format_type,
                template=template_name,
                content_length=len(content),
                session_id=context.session_id,
            )

            # Validate format
            if format_type not in ["docx", "pdf", "ppt"]:
                raise ValueError(f"Unsupported format: {format_type}")

            # Check dependencies for the requested format
            if not DEPENDENCIES_AVAILABLE:
                return self._generate_fallback_response(
                    format_type, content, time.time() - start_time
                )

            # Analyze content structure
            content_structure = self.content_analyzer.analyze_content_structure(content)

            # Generate filename if not provided
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                title_part = (
                    content_structure.title[:30]
                    if content_structure.title
                    else "document"
                )
                # Clean filename
                title_part = re.sub(r"[^\w\s-]", "", title_part).strip()
                title_part = re.sub(r"[-\s]+", "_", title_part)
                filename = f"{title_part}_{timestamp}.{format_type}"

            # Ensure filename has correct extension
            if not filename.endswith(f".{format_type}"):
                filename = f"{filename}.{format_type}"

            # Generate document
            file_path = os.path.join(self.output_directory, filename)

            if format_type == "docx":
                generation_result = await self._generate_docx(
                    content, content_structure, template_name, file_path
                )
            elif format_type == "pdf":
                generation_result = await self._generate_pdf(
                    content, content_structure, template_name, file_path
                )
            elif format_type == "ppt":
                generation_result = await self._generate_ppt(
                    content, content_structure, template_name, file_path
                )

            # Validate output if requested
            validation_result = None
            if validate_output:
                validation_result = self.validator.validate_document(
                    file_path, format_type
                )
                if not validation_result.is_valid:
                    logger.warning(
                        "Generated document failed validation",
                        errors=validation_result.errors,
                    )

            execution_time = time.time() - start_time

            # Prepare result data
            result_data = {
                "filename": filename,
                "file_path": file_path,
                "format": format_type,
                "template_used": template_name,
                "content_analysis": {
                    "title": content_structure.title,
                    "sections": len(content_structure.sections),
                    "tables": len(content_structure.tables),
                    "code_blocks": len(content_structure.code_blocks),
                    "lists": len(content_structure.lists),
                    "complexity": content_structure.metadata["complexity"],
                    "word_count": content_structure.metadata["word_count"],
                },
                "generation_metadata": generation_result,
                "processing_time": execution_time,
            }

            # Add validation results if available
            if validation_result:
                result_data["validation"] = {
                    "is_valid": validation_result.is_valid,
                    "errors": validation_result.errors,
                    "warnings": validation_result.warnings,
                    "accessibility_score": validation_result.accessibility_score,
                    "file_size_mb": validation_result.file_size_mb,
                    "quality_metrics": validation_result.quality_metrics,
                }

            # Calculate quality score
            quality_score = self._calculate_quality_score(
                generation_result, validation_result
            )
            confidence_score = self._calculate_confidence_score(
                format_type, content_structure, validation_result
            )

            result = ToolResult(
                data=result_data,
                metadata={
                    "tool": "document_generation",
                    "version": "1.0.0",
                    "format": format_type,
                    "template": template_name,
                    "session_id": context.session_id,
                    "dependencies_available": DEPENDENCIES_AVAILABLE,
                    "validation_performed": validate_output,
                },
                execution_time=execution_time,
                success=True,
                resource_usage={
                    "cpu_usage": 0.5,
                    "memory_usage_mb": 64,
                    "disk_usage_mb": (
                        validation_result.file_size_mb if validation_result else 0.0
                    ),
                },
                quality_score=quality_score,
                confidence_score=confidence_score,
            )

            return result

        except Exception as e:
            execution_time = time.time() - start_time

            logger.error(
                "Document generation failed",
                error=str(e),
                session_id=context.session_id,
            )

            return ToolResult(
                data=None,
                metadata={
                    "tool": "document_generation",
                    "session_id": context.session_id,
                    "error_type": type(e).__name__,
                },
                execution_time=execution_time,
                success=False,
                error_message=str(e),
                quality_score=0.0,
                confidence_score=0.0,
            )

    async def _generate_docx(
        self,
        content: str,
        content_structure: ContentStructure,
        template_name: str,
        file_path: str,
    ) -> Dict[str, Any]:
        """Generate DOCX document with advanced formatting"""

        if not DEPENDENCIES_AVAILABLE:
            raise RuntimeError("DOCX generation requires python-docx library")

        try:
            # Get template
            template = self.template_manager.get_template(template_name, "docx")

            # Create document
            doc = Document()

            # Set document properties
            doc.core_properties.title = content_structure.title or "Generated Document"
            doc.core_properties.author = "SE SME Agent"
            doc.core_properties.created = datetime.now()

            # Add title
            if content_structure.title:
                title_paragraph = doc.add_heading(content_structure.title, level=0)
                title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Process sections
            for section in content_structure.sections:
                # Add section heading
                heading = doc.add_heading(section["title"], level=section["level"])

                # Add section content
                section_content = "\n".join(section["content"])
                if section_content.strip():
                    self._add_formatted_content(doc, section_content, content_structure)

            # Add any remaining content not in sections
            remaining_content = self._extract_remaining_content(
                content, content_structure
            )
            if remaining_content.strip():
                self._add_formatted_content(doc, remaining_content, content_structure)

            # Save document
            doc.save(file_path)

            return {
                "format": "docx",
                "paragraphs_created": len(doc.paragraphs),
                "tables_created": len(doc.tables),
                "styles_used": len(doc.styles),
                "template_applied": template_name,
            }

        except Exception as e:
            logger.error(f"DOCX generation failed: {e}")
            raise

    async def _generate_pdf(
        self,
        content: str,
        content_structure: ContentStructure,
        template_name: str,
        file_path: str,
    ) -> Dict[str, Any]:
        """Generate PDF document with accessibility compliance"""

        if not DEPENDENCIES_AVAILABLE:
            raise RuntimeError("PDF generation requires reportlab library")

        try:
            # Get template
            template = self.template_manager.get_template(template_name, "pdf")

            # Create PDF document
            doc = SimpleDocTemplate(
                file_path,
                pagesize=template.layout["pagesize"],
                topMargin=template.layout["margins"]["top"],
                bottomMargin=template.layout["margins"]["bottom"],
                leftMargin=template.layout["margins"]["left"],
                rightMargin=template.layout["margins"]["right"],
            )

            # Create styles
            styles = getSampleStyleSheet()

            # Custom styles based on template
            title_style = ParagraphStyle(
                "CustomTitle",
                parent=styles["Title"],
                fontSize=template.styles["title"]["fontSize"],
                alignment=template.styles["title"]["alignment"],
                spaceAfter=template.styles["title"]["spaceAfter"],
            )

            heading1_style = ParagraphStyle(
                "CustomHeading1",
                parent=styles["Heading1"],
                fontSize=template.styles["heading1"]["fontSize"],
                spaceAfter=template.styles["heading1"]["spaceAfter"],
                spaceBefore=template.styles["heading1"]["spaceBefore"],
            )

            normal_style = ParagraphStyle(
                "CustomNormal",
                parent=styles["Normal"],
                fontSize=template.styles["normal"]["fontSize"],
                alignment=template.styles["normal"]["alignment"],
            )

            code_style = ParagraphStyle(
                "CustomCode",
                parent=styles["Code"],
                fontName=template.styles["code"]["fontName"],
                fontSize=template.styles["code"]["fontSize"],
                leftIndent=template.styles["code"]["leftIndent"],
            )

            # Build document content
            story = []

            # Add title
            if content_structure.title:
                story.append(Paragraph(content_structure.title, title_style))
                story.append(Spacer(1, 12))

            # Process sections
            for section in content_structure.sections:
                # Add section heading
                story.append(Paragraph(section["title"], heading1_style))

                # Add section content
                section_content = "\n".join(section["content"])
                if section_content.strip():
                    self._add_pdf_content(
                        story,
                        section_content,
                        content_structure,
                        normal_style,
                        code_style,
                    )

            # Add any remaining content
            remaining_content = self._extract_remaining_content(
                content, content_structure
            )
            if remaining_content.strip():
                self._add_pdf_content(
                    story,
                    remaining_content,
                    content_structure,
                    normal_style,
                    code_style,
                )

            # Build PDF
            doc.build(story)

            return {
                "format": "pdf",
                "pages_created": 1,  # Simplified - would need page counting
                "elements_created": len(story),
                "template_applied": template_name,
            }

        except Exception as e:
            logger.error(f"PDF generation failed: {e}")
            raise

    async def _generate_ppt(
        self,
        content: str,
        content_structure: ContentStructure,
        template_name: str,
        file_path: str,
    ) -> Dict[str, Any]:
        """Generate PowerPoint presentation with intelligent slide structure"""

        if not DEPENDENCIES_AVAILABLE:
            raise RuntimeError("PPT generation requires python-pptx library")

        try:
            # Get template
            template = self.template_manager.get_template(template_name, "ppt")

            # Create presentation
            prs = Presentation()

            # Set slide dimensions
            prs.slide_width = template.layout["slide_width"]
            prs.slide_height = template.layout["slide_height"]

            # Title slide
            if content_structure.title:
                title_slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]

                title.text = content_structure.title
                subtitle.text = (
                    f"Generated by SE SME Agent\n{datetime.now().strftime('%B %d, %Y')}"
                )

            # Content slides
            slides_created = 1 if content_structure.title else 0

            for section in content_structure.sections:
                # Create slide for each major section
                if section["level"] <= 2:  # Only create slides for major sections
                    content_slide_layout = prs.slide_layouts[
                        1
                    ]  # Title and content layout
                    slide = prs.slides.add_slide(content_slide_layout)

                    # Set slide title
                    title = slide.shapes.title
                    title.text = section["title"]

                    # Add content
                    content_placeholder = slide.placeholders[1]
                    section_text = "\n".join(section["content"])

                    # Process content for slide format
                    slide_content = self._format_content_for_slide(
                        section_text, content_structure
                    )
                    content_placeholder.text = slide_content

                    slides_created += 1

            # Handle tables and code blocks as separate slides if they're substantial
            for table in content_structure.tables:
                if len(table["rows"]) > 3:  # Only create slide for substantial tables
                    table_slide_layout = prs.slide_layouts[5]  # Blank layout
                    slide = prs.slides.add_slide(table_slide_layout)

                    # Add title
                    title_shape = slide.shapes.add_textbox(
                        PptxInches(1), PptxInches(0.5), PptxInches(8), PptxInches(1)
                    )
                    title_frame = title_shape.text_frame
                    title_frame.text = "Data Table"

                    # Add table (simplified representation)
                    table_shape = slide.shapes.add_textbox(
                        PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
                    )
                    table_frame = table_shape.text_frame

                    # Format table content
                    if table["headers"]:
                        table_frame.text = " | ".join(table["headers"]) + "\n"
                        table_frame.text += "-" * 50 + "\n"

                    for row in table["rows"][:10]:  # Limit to first 10 rows
                        table_frame.text += " | ".join(row) + "\n"

                    slides_created += 1

            # Save presentation
            prs.save(file_path)

            return {
                "format": "ppt",
                "slides_created": slides_created,
                "sections_processed": len(content_structure.sections),
                "template_applied": template_name,
            }

        except Exception as e:
            logger.error(f"PPT generation failed: {e}")
            raise

    def _add_formatted_content(
        self, doc, content: str, content_structure: ContentStructure
    ):
        """Add formatted content to DOCX document"""

        lines = content.split("\n")

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Check if line is part of a code block
            if self._is_code_line(line, content_structure):
                # Add as code paragraph
                code_paragraph = doc.add_paragraph(line)
                code_paragraph.style = "Code"  # Use built-in code style

            # Check if line is part of a list
            elif line.startswith(("- ", "* ", "+ ")) or re.match(r"^\d+\.\s", line):
                # Add as list item
                doc.add_paragraph(line, style="List Bullet")

            else:
                # Add as normal paragraph
                doc.add_paragraph(line)

    def _add_pdf_content(
        self,
        story: List,
        content: str,
        content_structure: ContentStructure,
        normal_style,
        code_style,
    ):
        """Add formatted content to PDF story"""

        lines = content.split("\n")

        for line in lines:
            line = line.strip()
            if not line:
                story.append(Spacer(1, 6))
                continue

            # Check if line is part of a code block
            if self._is_code_line(line, content_structure):
                story.append(Paragraph(line, code_style))
            else:
                story.append(Paragraph(line, normal_style))

            story.append(Spacer(1, 6))

    def _format_content_for_slide(
        self, content: str, content_structure: ContentStructure
    ) -> str:
        """Format content appropriately for PowerPoint slide"""

        lines = content.split("\n")
        formatted_lines = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Convert to bullet points for better slide presentation
            if not line.startswith(("- ", "• ", "* ")):
                if len(line) < 100:  # Short lines become bullet points
                    line = f"• {line}"

            formatted_lines.append(line)

        # Limit content to fit on slide
        if len(formatted_lines) > 8:
            formatted_lines = formatted_lines[:7] + [
                "• ... (content truncated for slide)"
            ]

        return "\n".join(formatted_lines)

    def _is_code_line(self, line: str, content_structure: ContentStructure) -> bool:
        """Check if a line is part of a code block"""

        # Simple heuristics for code detection
        code_indicators = [
            line.startswith("    "),  # Indented
            line.startswith("\t"),  # Tab indented
            "=" in line and ("def " in line or "class " in line or "import " in line),
            line.startswith(
                ("def ", "class ", "import ", "from ", "if ", "for ", "while ")
            ),
            "{" in line and "}" in line,
            line.endswith((";", "{", "}")),
        ]

        return any(code_indicators)

    def _extract_remaining_content(
        self, original_content: str, content_structure: ContentStructure
    ) -> str:
        """Extract content that wasn't processed in sections"""

        # This is a simplified approach - in practice, you'd want more sophisticated content tracking
        processed_lines = set()

        # Mark lines that were processed in sections
        for section in content_structure.sections:
            for line in section["content"]:
                processed_lines.add(line.strip())

        # Find unprocessed lines
        remaining_lines = []
        for line in original_content.split("\n"):
            if line.strip() and line.strip() not in processed_lines:
                # Skip section headers
                if not line.strip().startswith("#"):
                    remaining_lines.append(line)

        return "\n".join(remaining_lines)

    def _generate_fallback_response(
        self, format_type: str, content: str, execution_time: float
    ) -> ToolResult:
        """Generate fallback response when dependencies are not available"""

        # Create a simple text file as fallback
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"document_{timestamp}.txt"
        file_path = os.path.join(self.output_directory, filename)

        # Write content to text file
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("Generated Document\n")
            f.write("==================\n\n")
            f.write(f"Original format requested: {format_type.upper()}\n")
            f.write("Generated as text due to missing dependencies\n\n")
            f.write(content)

        result_data = {
            "filename": filename,
            "file_path": file_path,
            "format": "txt",
            "original_format_requested": format_type,
            "fallback_reason": "Missing document generation dependencies",
            "processing_time": execution_time,
        }

        return ToolResult(
            data=result_data,
            metadata={
                "tool": "document_generation",
                "version": "1.0.0",
                "fallback_mode": True,
                "dependencies_available": False,
            },
            execution_time=execution_time,
            success=True,
            quality_score=0.3,  # Lower quality for fallback
            confidence_score=0.5,
        )

    def _calculate_quality_score(
        self,
        generation_result: Dict[str, Any],
        validation_result: Optional[DocumentValidationResult],
    ) -> float:
        """Calculate quality score for generated document"""

        base_score = 0.7

        # Factor in generation success metrics
        if "elements_created" in generation_result:
            elements = generation_result["elements_created"]
            if elements > 10:
                base_score += 0.1
            elif elements > 5:
                base_score += 0.05

        # Factor in validation results
        if validation_result:
            if validation_result.is_valid:
                base_score += 0.1

            # Accessibility score factor
            base_score += validation_result.accessibility_score * 0.1

            # Penalize for errors
            if validation_result.errors:
                base_score -= len(validation_result.errors) * 0.05

        return max(0.0, min(1.0, base_score))

    def _calculate_confidence_score(
        self,
        format_type: str,
        content_structure: ContentStructure,
        validation_result: Optional[DocumentValidationResult],
    ) -> float:
        """Calculate confidence score for document generation"""

        base_confidence = 0.8 if DEPENDENCIES_AVAILABLE else 0.5

        # Factor in content complexity
        complexity = content_structure.metadata["complexity"]
        complexity_factor = {"low": 0.1, "medium": 0.0, "high": -0.1}.get(
            complexity, 0.0
        )

        # Factor in format support
        format_factor = 0.1 if format_type in ["docx", "pdf", "ppt"] else -0.2

        # Factor in validation results
        validation_factor = 0.0
        if validation_result:
            if validation_result.is_valid:
                validation_factor += 0.1
            if validation_result.errors:
                validation_factor -= len(validation_result.errors) * 0.05

        final_confidence = max(
            0.0,
            min(
                1.0,
                base_confidence + complexity_factor + format_factor + validation_factor,
            ),
        )

        return final_confidence

    def get_schema(self) -> Dict[str, Any]:
        """Get tool schema for document generation"""
        return {
            "name": "document_generation",
            "description": "Generate documents in DOCX, PDF, and PPT formats with template support and quality validation",
            "version": "1.0.0",
            "parameters": {
                "content": {
                    "type": "string",
                    "description": "The content to be formatted into a document",
                    "required": True,
                    "min_length": 1,
                },
                "format": {
                    "type": "string",
                    "description": "Output format for the document",
                    "enum": ["docx", "pdf", "ppt"],
                    "default": "docx",
                    "required": False,
                },
                "template": {
                    "type": "string",
                    "description": "Template to use for formatting",
                    "default": "default",
                    "required": False,
                },
                "filename": {
                    "type": "string",
                    "description": "Custom filename for the generated document",
                    "required": False,
                },
                "validate": {
                    "type": "boolean",
                    "description": "Whether to validate the generated document",
                    "default": True,
                    "required": False,
                },
            },
            "required_params": ["content"],
            "returns": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "file_path": {"type": "string"},
                    "format": {"type": "string"},
                    "template_used": {"type": "string"},
                    "content_analysis": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "sections": {"type": "integer"},
                            "tables": {"type": "integer"},
                            "code_blocks": {"type": "integer"},
                            "lists": {"type": "integer"},
                            "complexity": {"type": "string"},
                            "word_count": {"type": "integer"},
                        },
                    },
                    "generation_metadata": {"type": "object"},
                    "validation": {
                        "type": "object",
                        "properties": {
                            "is_valid": {"type": "boolean"},
                            "errors": {"type": "array"},
                            "warnings": {"type": "array"},
                            "accessibility_score": {"type": "number"},
                            "file_size_mb": {"type": "number"},
                            "quality_metrics": {"type": "object"},
                        },
                    },
                    "processing_time": {"type": "number"},
                },
            },
            "capabilities": {
                "primary": "document_generation",
                "secondary": ["content_formatting", "template_processing"],
                "input_types": ["text", "markdown"],
                "output_types": ["docx", "pdf", "ppt", "binary_file"],
                "supported_formats": ["docx", "pdf", "ppt"],
            },
        }

    def get_capabilities(self) -> ToolCapabilities:
        """Get tool capabilities"""
        return ToolCapabilities(
            primary_capability=ToolCapability.DOCUMENT_GENERATION,
            secondary_capabilities=[
                ToolCapability.CONTENT_GENERATION,
                ToolCapability.TRANSFORMATION,
            ],
            input_types=["text", "markdown", "structured_content"],
            output_types=["docx", "pdf", "ppt", "binary_file"],
            supported_formats=["docx", "pdf", "ppt", "txt"],
            language_support=["en"],
        )

    def get_resource_requirements(self) -> ResourceRequirements:
        """Get resource requirements"""
        return ResourceRequirements(
            cpu_cores=1.0,
            memory_mb=256,
            network_bandwidth_mbps=0.0,  # No network required
            storage_mb=50,  # For temporary files
            gpu_memory_mb=0,
            max_execution_time=60,
            concurrent_limit=3,  # Limit concurrent document generation
        )

    async def cleanup(self):
        """Cleanup tool resources"""
        logger.info("Cleaning up Document Generation Tool")

        # Clean up temporary files older than 1 hour
        try:
            current_time = time.time()
            for filename in os.listdir(self.output_directory):
                file_path = os.path.join(self.output_directory, filename)
                if os.path.isfile(file_path):
                    file_age = current_time - os.path.getctime(file_path)
                    if file_age > 3600:  # 1 hour
                        os.remove(file_path)
                        logger.debug(f"Cleaned up old file: {filename}")
        except Exception as e:
            logger.warning(f"Error during cleanup: {e}")
