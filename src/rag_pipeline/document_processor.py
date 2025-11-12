"""
Document Processor
Multi-format document processing and chunking
"""

from typing import Dict, Any, List, Optional
from dataclasses import dataclass
from pathlib import Path
import hashlib
from datetime import datetime
import asyncio
import mimetypes

# Document processing libraries
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation

from .models import (
    Chunk,
    DocumentMetadata,
    ProcessedDocument,
    DocumentType,
    ChunkingConfig,
)
from src.shared.logging import get_logger

logger = get_logger(__name__)

# ChunkingManager is imported lazily to avoid circular dependencies


# DocumentMetadata and Chunk are imported from models.py


@dataclass
class BatchProcessingResult:
    """Result of batch document processing"""

    total_documents: int
    successful: int
    failed: int
    processing_time: float
    results: List[Dict[str, Any]]
    errors: List[str] = None

    def __post_init__(self):
        if self.errors is None:
            self.errors = []


class DocumentProcessor:
    """Processes documents and creates chunks for embedding"""

    def __init__(self, chunking_config: ChunkingConfig = None):
        self.supported_formats = [".pdf", ".docx", ".pptx", ".txt", ".md"]

        # Initialize chunking system
        self.chunking_config = chunking_config or ChunkingConfig(
            strategy="hierarchical"
        )
        # Lazy import to avoid circular dependency
        from .chunking_strategies import ChunkingManager

        self.chunking_manager = ChunkingManager(self.chunking_config)

        # Legacy chunk sizes for backward compatibility
        self.chunk_sizes = self.chunking_config.chunk_size

        self._initialized = False
        self._processed_documents = {}  # Cache for deduplication
        self._processing_stats = {
            "total_processed": 0,
            "successful": 0,
            "failed": 0,
            "duplicates_skipped": 0,
        }

    async def initialize(self):
        """Initialize document processor"""
        if self._initialized:
            return

        logger.info("Initializing Document Processor")

        try:
            # Initialize document parsing libraries
            # Verify required libraries are available
            import PyPDF2
            import docx
            import pptx

            logger.info("Document parsing libraries verified")

            # Initialize chunking manager with lazy import
            from .chunking_strategies import ChunkingManager

            if not hasattr(self, "chunking_manager") or self.chunking_manager is None:
                self.chunking_manager = ChunkingManager(self.chunking_config)

            # Don't reinitialize _processing_stats - it's already set in constructor
            # Just add additional stats if needed
            self._processing_stats.update(
                {
                    "total_chunks_created": 0,
                    "processing_errors": 0,
                    "last_processed": None,
                }
            )

            logger.info("Document parsing libraries and chunking manager initialized")

        except ImportError as e:
            logger.warning(f"Some document parsing libraries not available: {e}")
            # Continue initialization - we can still process some document types
        except Exception as e:
            logger.error(f"Error initializing document processor: {e}")
            raise

        self._initialized = True
        logger.info("Document Processor initialized")

    async def process_document(
        self, file_path: str, metadata: Dict[str, Any] = None
    ) -> ProcessedDocument:
        """Process a document and create chunks with comprehensive error handling"""

        start_time = datetime.utcnow()
        processing_errors = []

        try:
            logger.info("Processing document", file_path=file_path)

            # Validate file exists
            if not Path(file_path).exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            # Detect document format and MIME type
            doc_format, mime_type = self._detect_format_and_mime(file_path)

            # Check for duplicates
            content_hash = await self._calculate_file_hash(file_path)
            if await self._is_duplicate(content_hash):
                logger.info(
                    "Duplicate document detected, skipping",
                    file_path=file_path,
                    content_hash=content_hash,
                )
                self._processing_stats["duplicates_skipped"] += 1
                return self._processed_documents[content_hash]

            # Extract content and metadata
            content, extracted_metadata = await self._extract_content_and_metadata(
                file_path, doc_format
            )

            # Merge metadata
            combined_metadata = self._merge_metadata(
                metadata or {}, extracted_metadata, file_path, mime_type
            )

            # Generate document ID
            doc_id = self._generate_document_id(file_path, content_hash)

            # Create chunks using enhanced chunking system
            chunks = await self.chunking_manager.create_chunks(
                content, doc_id, combined_metadata
            )

            # Calculate processing time
            processing_time = (datetime.utcnow() - start_time).total_seconds()

            processed_doc = ProcessedDocument(
                document_id=doc_id,
                original_path=file_path,
                content=content,
                chunks=chunks,
                metadata=combined_metadata,
                processing_time=processing_time,
                content_hash=content_hash,
                total_chunks=len(chunks),
                total_characters=len(content),
                processing_errors=processing_errors,
            )

            # Cache for deduplication
            self._processed_documents[content_hash] = processed_doc
            self._processing_stats["successful"] += 1
            self._processing_stats["total_processed"] += 1

            logger.info(
                "Document processed successfully",
                file_path=file_path,
                doc_id=doc_id,
                chunk_count=len(chunks),
                processing_time=processing_time,
            )

            return processed_doc

        except Exception as e:
            error_msg = f"Failed to process document {file_path}: {str(e)}"
            processing_errors.append(error_msg)
            logger.error(
                "Document processing failed", file_path=file_path, error=str(e)
            )

            self._processing_stats["failed"] += 1
            self._processing_stats["total_processed"] += 1

            # Return minimal processed document with error
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            return ProcessedDocument(
                document_id=f"error_{hashlib.md5(file_path.encode()).hexdigest()}",
                original_path=file_path,
                content="",
                chunks=[],
                metadata=DocumentMetadata(
                    source_path=file_path,
                    file_name=Path(file_path).name,
                    title=f"Error processing {Path(file_path).name}",
                    document_type=DocumentType.UNKNOWN,
                ),
                processing_time=processing_time,
                content_hash="",
                total_chunks=0,
                total_characters=0,
                processing_errors=processing_errors,
            )

    def _detect_format_and_mime(self, file_path: str) -> tuple[str, str]:
        """Detect document format from file extension and MIME type"""

        path = Path(file_path)
        extension = path.suffix.lower()

        # Get MIME type
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = "application/octet-stream"

        # Validate format support
        if extension in self.supported_formats:
            format_type = extension[1:]  # Remove the dot
        else:
            logger.warning(
                "Unsupported format detected", file_path=file_path, extension=extension
            )
            format_type = "unknown"

        return format_type, mime_type

    async def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA-256 hash of file for deduplication"""

        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception as e:
            logger.error(
                "Failed to calculate file hash", file_path=file_path, error=str(e)
            )
            # Fallback to path-based hash
            return hashlib.sha256(file_path.encode()).hexdigest()

    async def _is_duplicate(self, content_hash: str) -> bool:
        """Check if document is a duplicate based on content hash"""
        return content_hash in self._processed_documents

    async def _extract_content_and_metadata(
        self, file_path: str, doc_format: str
    ) -> tuple[str, DocumentMetadata]:
        """Extract text content and metadata from document"""

        try:
            if doc_format == "txt":
                return await self._extract_txt_with_metadata(file_path)
            elif doc_format == "md":
                return await self._extract_markdown_with_metadata(file_path)
            elif doc_format == "pdf":
                return await self._extract_pdf_with_metadata(file_path)
            elif doc_format == "docx":
                return await self._extract_docx_with_metadata(file_path)
            elif doc_format == "pptx":
                return await self._extract_pptx_with_metadata(file_path)
            else:
                # Fallback to text extraction
                content, _ = await self._extract_txt_with_metadata(file_path)
                return content, DocumentMetadata(
                    source_path=file_path,
                    file_name=Path(file_path).name,
                    document_type=DocumentType.UNKNOWN,
                )

        except Exception as e:
            logger.error("Content extraction failed", file_path=file_path, error=str(e))
            error_content = f"Error extracting content from {file_path}: {str(e)}"
            return error_content, DocumentMetadata(
                source_path=file_path,
                file_name=Path(file_path).name,
                title=f"Error: {Path(file_path).name}",
                document_type=DocumentType.UNKNOWN,
            )

    async def _extract_txt_with_metadata(
        self, file_path: str
    ) -> tuple[str, DocumentMetadata]:
        """Extract content and metadata from text file"""
        try:
            # Try UTF-8 first
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            encoding = "utf-8"
        except UnicodeDecodeError:
            # Fallback to latin-1
            with open(file_path, "r", encoding="latin-1") as f:
                content = f.read()
            encoding = "latin-1"

        # Get file stats
        path = Path(file_path)
        stat = path.stat()

        # Create metadata
        metadata = DocumentMetadata(
            source_path=file_path,
            file_name=path.name,
            file_size=stat.st_size,
            document_type=DocumentType.TXT,
            created_at=datetime.fromtimestamp(stat.st_ctime),
            modified_at=datetime.fromtimestamp(stat.st_mtime),
            title=path.stem,
            word_count=len(content.split()),
            mime_type="text/plain",
            encoding=encoding,
            category="text",
        )

        return content, metadata

    async def _extract_markdown_with_metadata(
        self, file_path: str
    ) -> tuple[str, DocumentMetadata]:
        """Extract content and metadata from markdown file"""
        try:
            # Try UTF-8 first
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            encoding = "utf-8"
        except UnicodeDecodeError:
            # Fallback to latin-1
            with open(file_path, "r", encoding="latin-1") as f:
                content = f.read()
            encoding = "latin-1"

        # Get file stats
        path = Path(file_path)
        stat = path.stat()

        # Extract title from markdown (first # header)
        title = path.stem
        lines = content.split("\n")
        for line in lines:
            if line.strip().startswith("# "):
                title = line.strip()[2:].strip()
                break

        # Extract tags from markdown (look for tags: or keywords: in front matter)
        tags = []
        keywords = []
        if content.startswith("---"):
            # Parse front matter
            try:
                front_matter_end = content.find("---", 3)
                if front_matter_end > 0:
                    front_matter = content[3:front_matter_end]
                    for line in front_matter.split("\n"):
                        if line.strip().startswith("tags:"):
                            tags_str = line.split(":", 1)[1].strip()
                            tags = [tag.strip() for tag in tags_str.split(",")]
                        elif line.strip().startswith("keywords:"):
                            keywords_str = line.split(":", 1)[1].strip()
                            keywords = [kw.strip() for kw in keywords_str.split(",")]
            except:
                pass  # Ignore front matter parsing errors

        # Create metadata
        metadata = DocumentMetadata(
            source_path=file_path,
            file_name=path.name,
            file_size=stat.st_size,
            document_type=DocumentType.MD,
            created_at=datetime.fromtimestamp(stat.st_ctime),
            modified_at=datetime.fromtimestamp(stat.st_mtime),
            title=title,
            word_count=len(content.split()),
            mime_type="text/markdown",
            encoding=encoding,
            category="documentation",
            tags=tags,
            keywords=keywords,
        )

        return content, metadata

    async def _extract_pdf_with_metadata(
        self, file_path: str
    ) -> tuple[str, DocumentMetadata]:
        """Extract content and metadata from PDF file"""
        try:
            content = ""
            metadata_dict = {}

            with open(file_path, "rb") as file:
                pdf_reader = pypdf.PdfReader(file)

                # Extract text from all pages
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                    except Exception as e:
                        logger.warning(
                            f"Failed to extract text from page {page_num + 1}: {e}"
                        )
                        continue

                # Extract metadata
                if pdf_reader.metadata:
                    metadata_dict = {
                        "title": pdf_reader.metadata.get("/Title", ""),
                        "author": pdf_reader.metadata.get("/Author", ""),
                        "subject": pdf_reader.metadata.get("/Subject", ""),
                        "creator": pdf_reader.metadata.get("/Creator", ""),
                        "producer": pdf_reader.metadata.get("/Producer", ""),
                        "creation_date": pdf_reader.metadata.get("/CreationDate", ""),
                        "modification_date": pdf_reader.metadata.get("/ModDate", ""),
                    }

            # Get file stats
            path = Path(file_path)
            stat = path.stat()

            # Parse dates if available
            creation_date = None
            modification_date = None
            try:
                if metadata_dict.get("creation_date"):
                    # PDF dates are in format D:YYYYMMDDHHmmSSOHH'mm'
                    date_str = metadata_dict["creation_date"].replace("D:", "")[:14]
                    creation_date = datetime.strptime(date_str, "%Y%m%d%H%M%S")
            except:
                creation_date = datetime.fromtimestamp(stat.st_ctime)

            try:
                if metadata_dict.get("modification_date"):
                    date_str = metadata_dict["modification_date"].replace("D:", "")[:14]
                    modification_date = datetime.strptime(date_str, "%Y%m%d%H%M%S")
            except:
                modification_date = datetime.fromtimestamp(stat.st_mtime)

            # Create metadata
            metadata = DocumentMetadata(
                source_path=file_path,
                file_name=path.name,
                file_size=stat.st_size,
                document_type=DocumentType.PDF,
                created_at=creation_date,
                modified_at=modification_date,
                title=metadata_dict.get("title") or path.stem,
                author=metadata_dict.get("author"),
                page_count=len(pdf_reader.pages),
                word_count=len(content.split()),
                mime_type="application/pdf",
                subject=metadata_dict.get("subject"),
                category="document",
            )

            return content, metadata

        except Exception as e:
            logger.error(f"Failed to extract PDF content: {e}")
            raise

    async def _extract_docx_with_metadata(
        self, file_path: str
    ) -> tuple[str, DocumentMetadata]:
        """Extract content and metadata from DOCX file"""
        try:
            doc = DocxDocument(file_path)

            # Extract text content
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    content += " | ".join(row_text) + "\n"

            # Get file stats
            path = Path(file_path)
            stat = path.stat()

            # Extract metadata from document properties
            props = doc.core_properties

            # Create metadata
            metadata = DocumentMetadata(
                source_path=file_path,
                file_name=path.name,
                file_size=stat.st_size,
                document_type=DocumentType.DOCX,
                created_at=props.created or datetime.fromtimestamp(stat.st_ctime),
                modified_at=props.modified or datetime.fromtimestamp(stat.st_mtime),
                title=props.title or path.stem,
                author=props.author,
                word_count=len(content.split()),
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                subject=props.subject,
                keywords=props.keywords.split(",") if props.keywords else [],
                category="document",
            )

            return content, metadata

        except Exception as e:
            logger.error(f"Failed to extract DOCX content: {e}")
            raise

    async def _extract_pptx_with_metadata(
        self, file_path: str
    ) -> tuple[str, DocumentMetadata]:
        """Extract content and metadata from PPTX file"""
        try:
            prs = Presentation(file_path)

            # Extract text content from all slides
            content = ""
            for slide_num, slide in enumerate(prs.slides):
                content += f"\n--- Slide {slide_num + 1} ---\n"

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"

                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            content += " | ".join(row_text) + "\n"

            # Get file stats
            path = Path(file_path)
            stat = path.stat()

            # Extract metadata from presentation properties
            props = prs.core_properties

            # Create metadata
            metadata = DocumentMetadata(
                source_path=file_path,
                file_name=path.name,
                file_size=stat.st_size,
                document_type=DocumentType.PPTX,
                created_at=props.created or datetime.fromtimestamp(stat.st_ctime),
                modified_at=props.modified or datetime.fromtimestamp(stat.st_mtime),
                title=props.title or path.stem,
                author=props.author,
                page_count=len(prs.slides),
                word_count=len(content.split()),
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                subject=props.subject,
                keywords=props.keywords.split(",") if props.keywords else [],
                category="presentation",
            )

            return content, metadata

        except Exception as e:
            logger.error(f"Failed to extract PPTX content: {e}")
            raise

    def _merge_metadata(
        self,
        user_metadata: Dict[str, Any],
        extracted_metadata: DocumentMetadata,
        file_path: str,
        mime_type: str,
    ) -> DocumentMetadata:
        """Merge user-provided metadata with extracted metadata"""

        # Start with extracted metadata
        merged = extracted_metadata

        # Override with user-provided metadata where available
        if user_metadata:
            for key, value in user_metadata.items():
                if hasattr(merged, key) and value is not None:
                    setattr(merged, key, value)

        # Ensure essential fields are set
        if not merged.mime_type:
            merged.mime_type = mime_type

        if not merged.title:
            merged.title = Path(file_path).stem

        # Add processing timestamp
        if not merged.modified_at:
            merged.modified_at = datetime.utcnow()

        return merged

    def _generate_document_id(self, file_path: str, content_hash: str) -> str:
        """Generate unique document ID"""

        # Create ID from file path and content hash
        path_hash = hashlib.md5(file_path.encode()).hexdigest()[:8]
        content_hash_short = content_hash[:8]

        return f"doc_{path_hash}_{content_hash_short}"

    def update_chunking_config(self, config: ChunkingConfig):
        """Update chunking configuration"""
        self.chunking_config = config
        # Lazy import to avoid circular dependency
        from .chunking_strategies import ChunkingManager

        self.chunking_manager = ChunkingManager(config)
        self.chunk_sizes = config.chunk_sizes
        logger.info("Updated chunking configuration", strategy=config.strategy)

    async def process_batch(
        self,
        file_paths: List[str],
        metadata_list: Optional[List[Dict[str, Any]]] = None,
        max_concurrent: int = 5,
    ) -> BatchProcessingResult:
        """Process multiple documents concurrently with comprehensive error handling"""

        start_time = datetime.utcnow()
        results = []
        errors = []

        # Ensure metadata list matches file paths
        if metadata_list is None:
            metadata_list = [{}] * len(file_paths)
        elif len(metadata_list) != len(file_paths):
            metadata_list.extend([{}] * (len(file_paths) - len(metadata_list)))

        logger.info(f"Starting batch processing of {len(file_paths)} documents")

        # Create semaphore to limit concurrent processing
        semaphore = asyncio.Semaphore(max_concurrent)

        async def process_single_with_semaphore(
            file_path: str, metadata: Dict[str, Any]
        ):
            async with semaphore:
                try:
                    result = await self.process_document(file_path, metadata)
                    return {
                        "file_path": file_path,
                        "status": "success",
                        "document_id": result.document_id,
                        "chunk_count": len(result.chunks),
                        "processing_time": result.processing_time,
                        "content_hash": result.content_hash,
                        "errors": result.processing_errors,
                    }
                except Exception as e:
                    error_msg = f"Failed to process {file_path}: {str(e)}"
                    errors.append(error_msg)
                    logger.error(
                        "Batch processing error", file_path=file_path, error=str(e)
                    )
                    return {
                        "file_path": file_path,
                        "status": "failed",
                        "error": str(e),
                        "processing_time": 0,
                    }

        # Process all documents concurrently
        tasks = [
            process_single_with_semaphore(file_path, metadata)
            for file_path, metadata in zip(file_paths, metadata_list)
        ]

        results = await asyncio.gather(*tasks, return_exceptions=True)

        # Process results and count successes/failures
        successful = 0
        failed = 0
        processed_results = []

        for result in results:
            if isinstance(result, Exception):
                failed += 1
                errors.append(f"Unexpected error: {str(result)}")
                processed_results.append({"status": "failed", "error": str(result)})
            elif result["status"] == "success":
                successful += 1
                processed_results.append(result)
            else:
                failed += 1
                processed_results.append(result)

        # Calculate total processing time
        total_time = (datetime.utcnow() - start_time).total_seconds()

        batch_result = BatchProcessingResult(
            total_documents=len(file_paths),
            successful=successful,
            failed=failed,
            processing_time=total_time,
            results=processed_results,
            errors=errors,
        )

        logger.info(
            f"Batch processing completed: {successful} successful, {failed} failed, {total_time:.2f}s total"
        )

        return batch_result

    async def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        try:
            return {
                "total_processed": self._processing_stats.get("total_processed", 0),
                "successful": self._processing_stats.get("successful", 0),
                "failed": self._processing_stats.get("failed", 0),
                "duplicates_skipped": self._processing_stats.get(
                    "duplicates_skipped", 0
                ),
                "success_rate": (
                    self._processing_stats.get("successful", 0)
                    / max(1, self._processing_stats.get("total_processed", 1))
                )
                * 100,
                "cached_documents": len(self._processed_documents),
                "supported_formats": self.supported_formats,
                "chunk_sizes": self.chunk_sizes,
                "total_chunks_created": self._processing_stats.get(
                    "total_chunks_created", 0
                ),
            }
        except Exception as e:
            logger.error(f"Error getting processing stats: {e}")
            return {
                "total_processed": 0,
                "successful": 0,
                "failed": 0,
                "duplicates_skipped": 0,
                "success_rate": 0.0,
                "cached_documents": 0,
                "supported_formats": self.supported_formats,
                "chunk_sizes": getattr(self, "chunk_sizes", 1000),
                "error": str(e),
            }

    async def clear_cache(self):
        """Clear the document processing cache"""
        self._processed_documents.clear()
        logger.info("Document processing cache cleared")

    async def health_check(self) -> Dict[str, Any]:
        """Check health of document processor"""
        stats = await self.get_processing_stats()
        return {
            "status": "healthy" if self._initialized else "not_initialized",
            "initialized": self._initialized,
            "processing_stats": stats,
        }

    async def shutdown(self):
        """Shutdown document processor"""
        logger.info("Shutting down Document Processor")
        # Cleanup resources if needed
